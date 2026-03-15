"""
Enhanced Document Processor with OCR and Multi-Format Support
Supports: Images (OCR), PDFs (with OCR fallback), Excel, PowerPoint, Word
Uses free tools: Tesseract OCR, pdfplumber, pandas, python-pptx
"""

import os
from typing import List, Tuple
from PIL import Image
import pytesseract
import pdfplumber
from pdf2image import convert_from_path
import pandas as pd
from pptx import Presentation
from docx import Document

# Configure Tesseract path (Windows)
# If Tesseract is not in PATH, set it here:
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def extract_text_from_image(image_path: str) -> str:
    """
    Extract text from image using Tesseract OCR
    Supports: PNG, JPG, JPEG, BMP, TIFF
    """
    try:
        image = Image.open(image_path)
        # Use English by default, can add more languages: lang='eng+hin+kan'
        text = pytesseract.image_to_string(image, lang='eng')
        return text.strip()
    except Exception as e:
        print(f"Image OCR Error: {e}")
        return ""

def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Extract text from PDF with OCR fallback for scanned PDFs
    First tries text extraction, if minimal text found, uses OCR
    """
    try:
        # Try text extraction first
        text = ""
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        
        # If very little text extracted, likely a scanned PDF - use OCR
        if len(text.strip()) < 100:
            print(f"PDF appears to be scanned, using OCR...")
            try:
                images = convert_from_path(pdf_path, dpi=300)
                text = ""
                for i, image in enumerate(images):
                    print(f"OCR processing page {i+1}/{len(images)}...")
                    text += pytesseract.image_to_string(image, lang='eng') + "\n"
            except Exception as ocr_error:
                print(f"PDF OCR Error: {ocr_error}")
                # Return whatever text we got from initial extraction
                pass
        
        return text.strip()
    except Exception as e:
        print(f"PDF Processing Error: {e}")
        return ""

def extract_text_from_excel(excel_path: str) -> str:
    """
    Extract text from Excel files
    Converts all sheets to text representation
    """
    try:
        # Read all sheets
        dfs = pd.read_excel(excel_path, sheet_name=None)
        text = ""
        
        for sheet_name, df in dfs.items():
            text += f"\n\n=== Sheet: {sheet_name} ===\n"
            # Convert DataFrame to string, removing index
            text += df.to_string(index=False, na_rep='')
        
        return text.strip()
    except Exception as e:
        print(f"Excel Processing Error: {e}")
        return ""

def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract text from PowerPoint files
    Extracts text from all slides and shapes
    """
    try:
        prs = Presentation(pptx_path)
        text = ""
        
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n\n=== Slide {i} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"PowerPoint Processing Error: {e}")
        return ""

def extract_text_from_docx(docx_path: str) -> str:
    """
    Extract text from Word documents
    Enhanced version with better formatting
    """
    try:
        doc = Document(docx_path)
        text = ""
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text += row_text + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"Word Processing Error: {e}")
        return ""

def smart_chunk(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """
    Smart chunking strategy that preserves semantic boundaries
    - Splits on paragraphs first
    - Then sentences if needed
    - Maintains context with overlap
    """
    if not text or len(text) == 0:
        return []
    
    chunks = []
    
    # Split into paragraphs first
    paragraphs = text.split('\n\n')
    
    current_chunk = ""
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        
        # If adding this paragraph exceeds chunk size
        if len(current_chunk) + len(para) > chunk_size:
            if current_chunk:
                chunks.append(current_chunk.strip())
                # Start new chunk with overlap
                words = current_chunk.split()
                overlap_text = " ".join(words[-overlap:]) if len(words) > overlap else current_chunk
                current_chunk = overlap_text + " " + para
            else:
                # Paragraph itself is too long, split by sentences
                sentences = para.split('. ')
                for sentence in sentences:
                    if len(current_chunk) + len(sentence) > chunk_size:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                            current_chunk = sentence
                        else:
                            # Even single sentence is too long, force split
                            chunks.append(sentence[:chunk_size])
                            current_chunk = sentence[chunk_size:]
                    else:
                        current_chunk += sentence + ". "
        else:
            current_chunk += para + "\n\n"
    
    # Add remaining chunk
    if current_chunk.strip():
        chunks.append(current_chunk.strip())
    
    return chunks

def process_document(file_path: str, file_type: str) -> Tuple[str, List[str]]:
    """
    Main document processor
    Returns: (full_text, chunks)
    """
    file_type = file_type.lower()
    text = ""
    
    # Route to appropriate processor
    if file_type in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
        text = extract_text_from_image(file_path)
    elif file_type == 'pdf':
        text = extract_text_from_pdf(file_path)
    elif file_type in ['xlsx', 'xls']:
        text = extract_text_from_excel(file_path)
    elif file_type == 'pptx':
        text = extract_text_from_pptx(file_path)
    elif file_type == 'docx':
        text = extract_text_from_docx(file_path)
    elif file_type == 'txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
    
    # Chunk the text
    chunks = smart_chunk(text, chunk_size=500, overlap=50)
    
    return text, chunks
